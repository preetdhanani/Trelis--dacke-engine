"""5.8 Renderer -- the fill engine (deterministic). Implements the seed-slide-
duplication mechanic (D3) with duplicate-before-fill ordering (D13): every
output slide is produced by duplicating a *pristine* seed slide from the
source template and filling the duplicate's named shapes -- never via
`add_slide(layout)` + `placeholders`, and never by re-cloning an
already-filled slide (confirmed necessary by the spike's negative control).
"""
import copy

from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from .models.chart import SUPPORTED_CHART_TYPES
from .qa_checker import check_overflow

IMAGE_RT = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"

# Chevron-flow diagram styling (5.7, M4). DEFAULT_CHEVRON_HEX is only the
# fallback for direct/unit-test calls that don't pass brand_colors -- the CLI
# always threads the tenant's real manifest brand colors through, so a second
# tenant's diagrams never silently render in tenant-1's blue.
DEFAULT_CHEVRON_HEX = "0957D1"
CHEVRON_GAP_IN = 0.12  # visual separation between adjacent chevrons
CHEVRON_HEIGHT_IN = 1.1  # fixed row height; short/wide for legible interior text
CHEVRON_FONT_PT = 12
CHEVRON_FONT_HEX = "FFFFFF"  # white text on a brand-colored chevron


def duplicate_slide(output_prs, seed_slide):
    """Duplicate `seed_slide` (from the pristine template) into `output_prs` as
    a new slide, with the seed's full content already in place. Filling still
    happens afterward via fill_slot().

    Deliberately does NOT call output_prs.slides.add_slide(): that convenience
    method calls `slide.shapes.clone_layout_placeholders(...)` internally,
    which touches `.shapes` -- a lazyproperty, cached on first access -- before
    the seed's content is swapped in below. Once poisoned, `.shapes` keeps
    returning that stale, empty collection for the rest of that Slide object's
    life, even after the cSld swap installs a real, populated shape tree.
    Going through the low-level part API instead avoids ever touching
    `.shapes` until after the swap is complete.
    """
    slide_layout = seed_slide.slide_layout
    rId, dest = output_prs.part.add_slide(slide_layout)
    output_prs.slides._sldIdLst.add_sldId(rId)

    src_cSld = seed_slide._element.find(qn("p:cSld"))
    dest_cSld = dest._element.find(qn("p:cSld"))
    new_cSld = copy.deepcopy(src_cSld)
    dest._element.replace(dest_cSld, new_cSld)

    # Remap every image relationship the clone carries (D13). A naive XML
    # clone copies r:embed ids verbatim; those ids are dangling in the new
    # slide part's own .rels unless explicitly re-related here -- this is the
    # exact failure the spike's negative control reproduced and confirmed.
    for blip in new_cSld.findall(".//" + qn("a:blip")):
        r_embed = blip.get(qn("r:embed"))
        if not r_embed:
            continue
        image_part = seed_slide.part.related_part(r_embed)
        new_rId = dest.part.relate_to(image_part, IMAGE_RT)
        blip.set(qn("r:embed"), new_rId)

    return dest


def find_shape(slide, name):
    for shp in slide.shapes:
        if shp.name == name:
            return shp
    return None


def _set_text_preserving_format(text_frame, value):
    paragraphs = text_frame.paragraphs
    p0 = paragraphs[0]
    if not p0.runs:
        p0.text = value
    else:
        p0.runs[0].text = value
        for extra_run in p0.runs[1:]:
            extra_run._r.getparent().remove(extra_run._r)
    for extra_p in paragraphs[1:]:
        extra_p._p.getparent().remove(extra_p._p)


def _set_bullets_preserving_format(text_frame, items):
    template_p = text_frame.paragraphs[0]._p
    parent = template_p.getparent()

    new_ps = []
    for item in items:
        p_clone = copy.deepcopy(template_p)
        r_elems = p_clone.findall(qn("a:r"))
        if r_elems:
            r_elems[0].find(qn("a:t")).text = item
            for extra in r_elems[1:]:
                p_clone.remove(extra)
        new_ps.append(p_clone)

    for p in list(text_frame.paragraphs):
        p._p.getparent().remove(p._p)
    for p_el in new_ps:
        parent.append(p_el)


def _remove_shape(shape):
    shape._element.getparent().remove(shape._element)


def _add_native_chart(slide, slot, chart_spec):
    """Insert a native (editable) python-pptx chart at the slot's geometry per
    D7/5.8-point-5 -- the *preferred* way to fill an exhibit slot, since it
    stays a real editable chart rather than a flat picture. The chart type name
    is resolved through SUPPORTED_CHART_TYPES so nothing outside the D7 matrix
    can ever reach add_chart().
    """
    xl_member = SUPPORTED_CHART_TYPES.get(chart_spec.chart_type)
    if xl_member is None:  # defensive: build_chart_spec should already guarantee this
        raise RuntimeError(f"unsupported chart type {chart_spec.chart_type!r} reached the renderer (D7)")
    xl_type = getattr(XL_CHART_TYPE, xl_member)

    data = chart_spec.data
    chart_data = CategoryChartData()
    chart_data.categories = data.categories
    for series_name, values in data.series.items():
        chart_data.add_series(series_name, values)

    geom = slot.geometry_in
    graphic_frame = slide.shapes.add_chart(
        xl_type,
        Inches(geom.left_in), Inches(geom.top_in),
        Inches(geom.width_in), Inches(geom.height_in),
        chart_data,
    )
    # keep the exhibit addressable by its manifest name (D3), same as the
    # picture path renames its inserted shape.
    graphic_frame.name = slot.shape_name

    chart = graphic_frame.chart
    if data.title:
        chart.has_title = True
        chart.chart_title.text_frame.text = data.title
    else:
        chart.has_title = False
    # a legend only earns its space when there's more than one series
    chart.has_legend = len(data.series) > 1
    return graphic_frame


def _add_native_diagram(slide, slot, diagram_data, brand_colors=None):
    """Insert a native (editable) chevron-flow diagram at the slot's geometry
    (5.7/M4): N MSO_SHAPE.CHEVRON autoshapes in a single left-to-right row,
    each shape's own arrow implying "next step" -- no connector/arrowhead API
    needed. The chevrons are grouped into ONE GroupShape renamed to
    slot.shape_name so the slot stays addressable by a single manifest name
    (D3), exactly as _add_native_chart renames its graphic frame; each chevron
    inside the group remains a real, individually editable autoshape --
    nothing here is flattened to a picture.
    """
    steps = diagram_data.steps
    n = len(steps)
    geom = slot.geometry_in

    chevron_height_in = min(CHEVRON_HEIGHT_IN, geom.height_in)
    total_gap_in = CHEVRON_GAP_IN * (n - 1)
    chevron_width_in = (geom.width_in - total_gap_in) / n
    top_in = geom.top_in + (geom.height_in - chevron_height_in) / 2
    # "primary" is the one tenant-agnostic color key every manifest's brand.colors
    # is expected to define for generated visual elements like this -- looking up
    # a semantically-named key (e.g. the reference tenant's own "blue_primary")
    # would silently fall back to DEFAULT_CHEVRON_HEX for any other tenant whose
    # palette doesn't happen to use that exact name (caught by M7's second tenant).
    hex_color = (brand_colors or {}).get("primary", DEFAULT_CHEVRON_HEX)

    chevrons = []
    for i, label in enumerate(steps):
        left_in = geom.left_in + i * (chevron_width_in + CHEVRON_GAP_IN)
        shp = slide.shapes.add_shape(
            MSO_SHAPE.CHEVRON,
            Inches(left_in), Inches(top_in),
            Inches(chevron_width_in), Inches(chevron_height_in),
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor.from_string(hex_color)
        shp.line.fill.background()  # flat, no outline
        tf = shp.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.text = label
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        for run in p0.runs:
            run.font.size = Pt(CHEVRON_FONT_PT)
            run.font.color.rgb = RGBColor.from_string(CHEVRON_FONT_HEX)
        chevrons.append(shp)

    group = slide.shapes.add_group_shape(chevrons)
    group.name = slot.shape_name  # keep the slot addressable as ONE shape (D3)
    return group


def _run_font_size_pt(text_frame):
    """Best-effort read of the first run's font size (points), for the QA
    overflow estimate (5.9/M5). None if the shape has no run or the size
    isn't explicitly set on it (inherited from master/theme) -- in that case
    the overflow check is skipped rather than guessing a size (D9)."""
    paragraphs = text_frame.paragraphs
    if not paragraphs or not paragraphs[0].runs:
        return None
    size = paragraphs[0].runs[0].font.size
    return size.pt if size is not None else None


def fill_slot(slide, slot, value, image_source_path=None, chart_spec=None,
              diagram_spec=None, brand_colors=None, layout=None, slide_height_in=None):
    """Returns (filled, overflow_reason): `filled` is True if the slot was
    actually filled, False if it was left as the seed's own placeholder
    content. Image slots with no source are skipped rather than raising: a
    slide that can't get a real exhibit must be flagged (needs_review), never
    silently broken or crashed on (D9) -- the caller (render_slide) surfaces
    which required slots were skipped so the CLI can flag them.

    `overflow_reason` is the QA/Overflow Checker's (5.9/M5) verdict for
    text/bullets slots -- None unless the estimated text height exceeds the
    slot's safe growth budget (see qa_checker.py). Only computed when both
    `layout` and `slide_height_in` are supplied (opt-in, so existing callers
    that don't care about overflow are unaffected); always None for image
    slots, which have no analogous check.

    An image-typed exhibit slot can be filled three ways (5.8 point 5, D7):
    a native `chart_spec` (preferred -- stays editable) takes precedence,
    then a native `diagram_spec` (chevron-flow, also editable, 5.7/M4), else
    a rendered `image_source_path` picture; with none of them, the grey
    placeholder rect is left in place.
    """
    shape = find_shape(slide, slot.shape_name)
    if shape is None:
        raise RuntimeError(f"named shape '{slot.shape_name}' not found on duplicated slide")

    if slot.type == "text":
        if value is None:
            return False, None
        text = value if not slot.max_chars or len(value) <= slot.max_chars else value[: slot.max_chars - 1].rstrip() + "…"
        _set_text_preserving_format(shape.text_frame, text)
        overflow_reason = None
        if layout is not None and slide_height_in is not None:
            font_size_pt = _run_font_size_pt(shape.text_frame)
            overflow_reason = check_overflow(text, font_size_pt, slot, layout, slide_height_in)
        return True, overflow_reason

    elif slot.type == "bullets":
        if not value:
            return False, None
        items = list(value)
        if slot.max_items:
            items = items[: slot.max_items]
        if slot.max_chars_per_item:
            items = [
                i if len(i) <= slot.max_chars_per_item else i[: slot.max_chars_per_item - 1].rstrip() + "…"
                for i in items
            ]
        _set_bullets_preserving_format(shape.text_frame, items)
        overflow_reason = None
        if layout is not None and slide_height_in is not None:
            font_size_pt = _run_font_size_pt(shape.text_frame)
            overflow_reason = check_overflow(items, font_size_pt, slot, layout, slide_height_in)
        return True, overflow_reason

    elif slot.type == "image":
        if chart_spec is not None:
            _add_native_chart(slide, slot, chart_spec)
            _remove_shape(shape)  # remove the grey placeholder rect, per fill_protocol
            return True, None
        if diagram_spec is not None:
            _add_native_diagram(slide, slot, diagram_spec, brand_colors=brand_colors)
            _remove_shape(shape)  # remove the grey placeholder rect, per fill_protocol
            return True, None
        if image_source_path is None:
            return False, None  # grey placeholder rect stays -- caller flags needs_review
        geom = slot.geometry_in
        picture = slide.shapes.add_picture(
            str(image_source_path),
            Inches(geom.left_in), Inches(geom.top_in),
            width=Inches(geom.width_in), height=Inches(geom.height_in),
        )
        # add_picture() names the new shape generically ("Picture N"); rename it
        # back to the slot's shape_name so manifest addressing (D3) still
        # resolves this slot on a future lookup (spike S1 finding).
        picture.name = slot.shape_name
        _remove_shape(shape)  # remove the grey placeholder rect, per fill_protocol
        return True, None

    else:
        raise RuntimeError(f"unknown slot type: {slot.type}")


def render_slide(output_prs, seed_slide, layout, spec, image_paths=None, chart_specs=None,
                 diagram_specs=None, brand_colors=None, slide_height_in=None):
    """Duplicate `seed_slide` into `output_prs` and fill it per `layout`'s
    slots from `spec.slots` (D3: duplicate first; D13: fill only the
    duplicate, never the seed).

    Image-typed exhibit slots are filled from `chart_specs[shape_name]` (a
    native chart, preferred), `diagram_specs[shape_name]` (a native
    chevron-flow diagram, 5.7/M4), or `image_paths[shape_name]` (a picture);
    a slot with none of them is left as the grey placeholder. `brand_colors`
    is the manifest's brand color dict, used to theme diagram shapes.
    `slide_height_in` (the manifest's slide_dimensions_in height) enables the
    QA/Overflow Checker (5.9/M5) for text/bullets slots; omit it to skip that
    check entirely (e.g. tests that don't care about overflow).

    Returns (new_slide, skipped_required_slots, overflow_warnings):
    `skipped_required_slots` lists the shape_name of every *required* slot
    left unfilled (only possible for image slots with no chart, diagram, or
    picture source -- see fill_slot); `overflow_warnings` lists
    (shape_name, reason) for text/bullets slots whose estimated height
    exceeds their safe growth budget. Both let the caller flag the slide
    needs_review instead of shipping it silently incomplete or broken (D9).
    """
    image_paths = image_paths or {}
    chart_specs = chart_specs or {}
    diagram_specs = diagram_specs or {}
    new_slide = duplicate_slide(output_prs, seed_slide)
    skipped_required = []
    overflow_warnings = []
    for slot in layout.slots:
        if slot.type == "image":
            filled, _ = fill_slot(
                new_slide, slot, None,
                image_source_path=image_paths.get(slot.shape_name),
                chart_spec=chart_specs.get(slot.shape_name),
                diagram_spec=diagram_specs.get(slot.shape_name),
                brand_colors=brand_colors,
            )
        else:
            filled, overflow_reason = fill_slot(
                new_slide, slot, spec.slots.get(slot.shape_name),
                layout=layout, slide_height_in=slide_height_in,
            )
            if overflow_reason:
                overflow_warnings.append((slot.shape_name, overflow_reason))
        if not filled and slot.required:
            skipped_required.append(slot.shape_name)
    return new_slide, skipped_required, overflow_warnings


def replace_rendered_slide(output_prs, position, seed_slide, layout, spec, **render_kwargs):
    """Re-render the slide currently at `position` in output_prs.slides (M6,
    5.10): re-duplicate the *pristine* seed slide fresh and fill it per the
    updated spec -- the exact same fill path as initial generation (D10/D11)
    -- never re-clone the already-rendered output slide (D13). The fresh
    slide is then swapped into `position` so slide order is preserved (FR-1)
    and the stale slide is discarded; every other slide is untouched (FR-6).

    `render_kwargs` passes through to render_slide unchanged (e.g.
    chart_specs/diagram_specs/brand_colors/slide_height_in) so a text-only
    revision doesn't lose that slide's existing chart/diagram.

    Returns (new_slide, skipped_required, overflow_warnings), the same
    contract as render_slide.
    """
    id_lst = output_prs.slides._sldIdLst
    old_sld_id = list(id_lst)[position]

    new_slide, skipped, overflow = render_slide(output_prs, seed_slide, layout, spec, **render_kwargs)

    new_sld_id = list(id_lst)[-1]  # render_slide always appends at the end
    id_lst.remove(new_sld_id)
    id_lst.insert(position, new_sld_id)
    id_lst.remove(old_sld_id)
    return new_slide, skipped, overflow


def strip_seed_slides(prs, count):
    """Remove the first `count` slides -- the pristine seed/index slides that
    came along when the output presentation was opened from a copy of the
    template -- leaving only the newly-rendered slides. Must be called *after*
    all render_slide() calls; removing by index while more slides are still
    being appended would shift indices under you.
    """
    id_lst = prs.slides._sldIdLst
    sld_ids = list(id_lst)
    for i in range(count - 1, -1, -1):
        id_lst.remove(sld_ids[i])
