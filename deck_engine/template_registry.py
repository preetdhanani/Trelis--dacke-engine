"""5.1 Template Registry & Manifest Loader.

Resolves a tenant's template .pptx + manifest and validates the manifest
against the real file: every layout's slide_index must exist in the file, and
every slot's shape_name must be a real shape.name on that slide. Fails fast at
load time -- manifest/template drift must never surface as a render-time
error (SDD 5.1's stated failure mode).
"""
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation

from .config import TenantConfig, get_tenant_config
from .models.manifest import Manifest


class ManifestValidationError(Exception):
    """Manifest and the real .pptx disagree, or either file is missing. Fatal at load."""


@dataclass
class TenantAssets:
    tenant_id: str
    manifest: Manifest
    template_path: Path

    def open_template(self):
        """Return a fresh Presentation loaded from the pristine template file.

        Called once per render run to get the source-of-seed-slides object.
        Every call re-reads from disk, so the result is always pristine --
        never a previously-rendered/filled presentation (D13).
        """
        return Presentation(str(self.template_path))


def _validate(manifest: Manifest, prs) -> list[str]:
    errors: list[str] = []
    n_slides = len(prs.slides)
    for layout in manifest.layouts:
        if layout.slide_index >= n_slides:
            errors.append(
                f"layout '{layout.layout_id}': slide_index={layout.slide_index} "
                f"out of range (template has {n_slides} slides)"
            )
            continue
        slide = prs.slides[layout.slide_index]
        shape_names = {shp.name for shp in slide.shapes}
        for slot in layout.slots:
            if slot.shape_name not in shape_names:
                errors.append(
                    f"layout '{layout.layout_id}' slide_index={layout.slide_index}: "
                    f"shape '{slot.shape_name}' not found on that slide "
                    f"(have: {sorted(shape_names)})"
                )
    return errors


def load_tenant_assets(tenant_id: str) -> TenantAssets:
    return load_tenant_assets_from_config(get_tenant_config(tenant_id))


def load_tenant_assets_from_config(config: TenantConfig) -> TenantAssets:
    if not config.template_path.exists():
        raise ManifestValidationError(f"template not found: {config.template_path}")
    if not config.manifest_path.exists():
        raise ManifestValidationError(f"manifest not found: {config.manifest_path}")

    manifest = Manifest.model_validate_json(config.manifest_path.read_text(encoding="utf-8"))
    prs = Presentation(str(config.template_path))

    errors = _validate(manifest, prs)
    if errors:
        raise ManifestValidationError(
            f"manifest/template mismatch for tenant '{config.tenant_id}':\n"
            + "\n".join(f"  - {e}" for e in errors)
        )

    return TenantAssets(tenant_id=config.tenant_id, manifest=manifest, template_path=config.template_path)
