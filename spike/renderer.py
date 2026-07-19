"""S1 -- throwaway renderer core: seed-slide duplication + named-shape fill.

Implements the mechanic described in SDD D3/D13/5.8:
  1. Duplicate the pristine seed slide (XML-clone), never a previously-filled slide.
  2. On the duplicate, find shapes by shape.name and fill per fill_protocol.
  3. Image slots: insert the real picture, then remove the grey placeholder rect.

Slide duplication here clones the seed slide's <p:cSld> (background + full shape
tree) into a slide part python-pptx has legitimately created via add_slide(), then
remaps every image relationship (r:embed) the clone carries so embedded pictures
survive the copy. This is the exact failure mode D13 warns about: a naive clone
that copies the XML but not the relationship copies a *dangling* r:embed id, so
the picture silently breaks. See s1_negative_control.py for a deliberate repro.
"""
import copy
import json
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

TEMPLATE_PATH = Path(__file__).parent.parent / "Templates" / "MasterDeck.pptx"
MANIFEST_PATH = Path(__file__).parent.parent / "Templates" / "MasterDeck_manifest.json"

IMAGE_RT = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"


def load_manifest():
    return json.loads(MANIFEST_PATH.read_text(encoding="utf-8"))


def duplicate_slide(prs, source_slide, remap_images=True):
    """Duplicate `source_slide` (must belong to a *possibly different* Presentation
    than `prs`, or the same one) into `prs` as a new slide. Returns the new slide.

    remap_images=False deliberately reproduces the D13 bug for the negative control.
    """
    # Deliberately NOT prs.slides.add_slide(): that convenience method calls
    # slide.shapes.clone_layout_placeholders(...) internally, which touches the
    # `.shapes` lazyproperty *before* we've swapped in the seed slide's content.
    # `.shapes` caches its SlideShapes wrapper on first access (bound to whatever
    # spTree existed then) -- once poisoned, every later `.shapes` access on this
    # same Slide object keeps returning that stale, now-detached, empty collection,
    # even after the cSld swap below installs a brand new, populated spTree. Going
    # through the low-level part API instead avoids ever touching `.shapes` until
    # after the swap is complete.
    slide_layout = source_slide.slide_layout
    rId, dest = prs.part.add_slide(slide_layout)
    prs.slides._sldIdLst.add_sldId(rId)

    src_cSld = source_slide._element.find(qn("p:cSld"))
    dest_cSld = dest._element.find(qn("p:cSld"))
    new_cSld = copy.deepcopy(src_cSld)
    dest._element.replace(dest_cSld, new_cSld)

    if remap_images:
        for blip in new_cSld.findall(".//" + qn("a:blip")):
            r_embed = blip.get(qn("r:embed"))
            if not r_embed:
                continue
            image_part = source_slide.part.related_part(r_embed)
            new_rId = dest.part.relate_to(image_part, IMAGE_RT)
            blip.set(qn("r:embed"), new_rId)
    # else: leave dangling r:embed ids pointing at relationship ids that don't
    # exist in dest's own .rels -- this is the bug D13 describes.

    return dest


def find_shape(slide, name):
    for shp in slide.shapes:
        if shp.name == name:
            return shp
    return None


def _set_text_preserving_format(text_frame, value):
    """Write `value` into the first run of the first paragraph, preserving that
    run's rPr (font/size/color) so styling is inherited rather than reset. Extra
    paragraphs/runs beyond the first are dropped.
    """
    paragraphs = text_frame.paragraphs
    p0 = paragraphs[0]
    if not p0.runs:
        p0.text = value
    else:
        p0.runs[0].text = value
        for extra_run in p0.runs[1:]:
            extra_run._r.getparent().remove(extra_run._r)
    for extra_p in paragraphs[1:]:
        extra_p._p.getparent().remove(extra_p._p)


def _set_bullets_preserving_format(text_frame, items, max_items=None, max_chars_per_item=None):
    items = list(items)
    if max_items:
        items = items[:max_items]
    if max_chars_per_item:
        items = [i if len(i) <= max_chars_per_item else i[: max_chars_per_item - 1].rstrip() + "…" for i in items]

    paragraphs = text_frame.paragraphs
    template_p = paragraphs[0]._p
    parent = template_p.getparent()

    # Build one <a:p> per item by cloning the template paragraph's formatting.
    new_ps = []
    for item in items:
        p_clone = copy.deepcopy(template_p)
        # find first run in the clone and set its text; drop other runs
        r_elems = p_clone.findall(qn("a:r"))
        if r_elems:
            t_el = r_elems[0].find(qn("a:t"))
            t_el.text = item
            for extra in r_elems[1:]:
                p_clone.remove(extra)
        else:
            # no runs at all (shouldn't happen on a real template paragraph) -- fallback
            pass
        new_ps.append(p_clone)

    # remove all original paragraphs, insert the new ones in order
    for p in list(text_frame.paragraphs):
        p._p.getparent().remove(p._p)
    for p_el in new_ps:
        parent.append(p_el)


def _remove_shape(shape):
    shape._element.getparent().remove(shape._element)


def fill_slot(slide, slot, spec_value, image_source_path=None):
    shape = find_shape(slide, slot["shape_name"])
    if shape is None:
        raise RuntimeError(f"named shape '{slot['shape_name']}' not found on duplicated slide")

    if slot["type"] == "text":
        if spec_value is None:
            return
        max_chars = slot.get("max_chars")
        value = spec_value if not max_chars or len(spec_value) <= max_chars else spec_value[: max_chars - 1].rstrip() + "…"
        _set_text_preserving_format(shape.text_frame, value)

    elif slot["type"] == "bullets":
        if not spec_value:
            return
        _set_bullets_preserving_format(
            shape.text_frame, spec_value,
            max_items=slot.get("max_items"), max_chars_per_item=slot.get("max_chars_per_item"),
        )

    elif slot["type"] == "image":
        geom = slot["geometry_in"]
        left, top = Inches(geom["left_in"]), Inches(geom["top_in"])
        width, height = Inches(geom["width_in"]), Inches(geom["height_in"])
        if image_source_path is None:
            raise RuntimeError(f"no image provided for slot '{slot['shape_name']}'")
        picture = slide.shapes.add_picture(str(image_source_path), left, top, width=width, height=height)
        # add_picture() names the new shape generically ("Picture N"); rename it back
        # to the slot's shape_name so the manifest's stable-name addressing (D3) still
        # resolves this slot on a future lookup (e.g. a field-edit regen of just this image).
        picture.name = slot["shape_name"]
        _remove_shape(shape)  # remove the grey placeholder rect, per fill_protocol step 6

    else:
        raise RuntimeError(f"unknown slot type: {slot['type']}")


def render_slide(output_prs, template_prs, layout, slide_spec, image_paths=None):
    """layout: one manifest layout dict. slide_spec: {shape_name: value}. image_paths:
    {shape_name: path} for image slots. Returns the new slide in output_prs."""
    image_paths = image_paths or {}
    seed_slide = template_prs.slides[layout["slide_index"]]
    new_slide = duplicate_slide(output_prs, seed_slide)
    for slot in layout["slots"]:
        name = slot["shape_name"]
        if slot["type"] == "image":
            fill_slot(new_slide, slot, None, image_source_path=image_paths.get(name))
        else:
            fill_slot(new_slide, slot, slide_spec.get(name))
    return new_slide
