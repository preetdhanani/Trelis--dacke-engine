"""S1 negative control -- empirically reproduce the D13 failure mode.

D13 says: never duplicate a slide that already has a picture inserted; only
ever duplicate a *pristine* seed slide, because XML-cloning a slide loses
embedded-image *relationships* unless they are explicitly remapped. This
script proves that claim directly, within a single Presentation/package
(mirroring how the real renderer always operates -- source and destination
are the same open package, never two independently-loaded files):

  1. Duplicate+fill the two_column seed slide once (normal renderer path) --
     this is the "already has a picture inserted" slide D13 warns about.
  2. Re-clone THAT already-filled slide a second time two ways:
       a. WITHOUT relationship remap (remap_images=False) -- reproduces the bug.
       b. WITH relationship remap (remap_images=True)  -- shows the mitigation works.
  3. Save one deck containing all three (already-filled, naive-reclone,
     fixed-reclone) and check whether each still resolves its picture.
"""
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

from renderer import TEMPLATE_PATH, load_manifest, duplicate_slide, render_slide

OUT_DIR = Path(__file__).parent / "output"
ASSETS = Path(__file__).parent / "assets"
OUT_PATH = OUT_DIR / "S1_negative_control.pptx"


def find_blip_rid(slide, shape_name):
    for shp in slide.shapes:
        if shp.name == shape_name:
            blip = shp._element.find(".//" + qn("a:blip"))
            return blip.get(qn("r:embed")) if blip is not None else None
    return None


def check_picture_resolves(slide, shape_name):
    try:
        shp = next(s for s in slide.shapes if s.name == shape_name)
    except StopIteration:
        return {"shape_found": False}
    rid = find_blip_rid(slide, shape_name)
    try:
        blob_len = len(shp.image.blob)
        return {"shape_found": True, "r_embed_id": rid, "resolved": True, "image_bytes": blob_len}
    except Exception as e:
        return {"shape_found": True, "r_embed_id": rid, "resolved": False, "error": f"{type(e).__name__}: {e}"}


def main():
    prs = Presentation(str(TEMPLATE_PATH))
    manifest = load_manifest()
    layout = next(l for l in manifest["layouts"] if l["layout_id"] == "two_column")

    # Step 1: normal render -> an "already filled" slide (has a real picture).
    spec = {"IM_TWOCOL_TITLE": "Negative control fixture", "IM_TWOCOL_BODY": ["fixture bullet one", "fixture bullet two"]}
    images = {"IM_TWOCOL_IMAGE": ASSETS / "two_col_support.jpg"}
    already_filled = render_slide(prs, prs, layout, spec, image_paths=images)

    # Step 2: re-clone that already-filled slide -- the operation D13 forbids.
    naive_reclone = duplicate_slide(prs, already_filled, remap_images=False)
    fixed_reclone = duplicate_slide(prs, already_filled, remap_images=True)

    # Drop the 9 original pristine seed slides; keep only our 3 test slides, in order.
    id_lst = prs.slides._sldIdLst
    for sld_id in list(id_lst)[:9]:
        id_lst.remove(sld_id)

    prs.save(str(OUT_PATH))
    print(f"wrote {OUT_PATH} with {len(prs.slides)} slides (already_filled, naive_reclone, fixed_reclone)")

    # Step 3: reopen FROM DISK (the only thing that matters -- does the saved file
    # still resolve the picture) and check all three.
    reopened = Presentation(str(OUT_PATH))
    labels = ["already_filled (original)", "naive_reclone (no remap)", "fixed_reclone (with remap)"]
    for label, slide in zip(labels, reopened.slides):
        print(f"\n{label}:")
        print(" ", check_picture_resolves(slide, "IM_TWOCOL_IMAGE"))


if __name__ == "__main__":
    main()
