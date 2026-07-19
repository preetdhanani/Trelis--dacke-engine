"""S1 -- full duplicate -> fill -> save round-trip against the real template.

For each of the 8 manifest layouts: duplicate its pristine seed slide out of
MasterDeck.pptx, fill every named slot with sample content (text,
bullets, and a real inserted picture for image slots), then save a single
output deck containing only the 8 filled slides (the original 9 seed/index
slides are dropped, mirroring how the real Assembler starts from a template
copy and builds up only the rendered slides -- SDD 5.11).
"""
import json
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

from renderer import TEMPLATE_PATH, load_manifest, render_slide

ASSETS = Path(__file__).parent / "assets"
OUT_DIR = Path(__file__).parent / "output"
OUT_DIR.mkdir(exist_ok=True)
OUT_PATH = OUT_DIR / "S1_filled_deck.pptx"

IMAGE_FOR_SLOT = {
    "IM_IMAGEONLY_HERO": ASSETS / "image_only_hero.jpg",
    "IM_TWOCOL_IMAGE": ASSETS / "two_col_support.jpg",
    "IM_EXHIBIT_VISUAL": ASSETS / "exhibit_chart.jpg",
}

SAMPLE_SPECS = {
    "title": {
        "IM_TITLE_DECKTITLE": "Q3 Market Entry Strategy: Southeast Asia Expansion",
        "IM_TITLE_SUBTITLE": "Prepared for the Executive Steering Committee",
        "IM_TITLE_CLAIM": "AI Built to Deliver.",
        "IM_TITLE_META": "Sample Consulting · July 2026 · Confidential",
    },
    "section_divider": {
        "IM_SECTION_NUMBER": "01",
        "IM_SECTION_TITLE": "Market Landscape & Opportunity",
        "IM_SECTION_KICKER": "Sizing the addressable market and competitive dynamics",
    },
    "text_bullets": {
        "IM_BULLETS_TITLE": "Three structural tailwinds support entry now",
        "IM_BULLETS_KICKER": "Demand-side signals across the region are converging",
        "IM_BULLETS_BODY": [
            "Regional GDP growth outpaces developed markets by 3.2x over the last five years",
            "Regulatory reform in 2025 removed the primary foreign-ownership barrier",
            "Two of three target segments remain under-served by incumbent players",
            "Digital payments infrastructure now covers 78% of the urban population",
            "Talent costs remain 40-55% below comparable developed-market benchmarks",
        ],
    },
    "image_only": {
        "IM_IMAGEONLY_CAPTION": "Figure 1 — Proposed regional distribution network, illustrative",
    },
    "two_column": {
        "IM_TWOCOL_TITLE": "A phased entry de-risks capital commitment",
        "IM_TWOCOL_BODY": [
            "Phase 1: pilot in a single metro to validate unit economics",
            "Phase 2: expand distribution once CAC payback is under 9 months",
            "Phase 3: full regional rollout with a local JV partner",
            "Capital exposure stays under $4M until Phase 2 gate criteria are met",
        ],
    },
    "exhibit_data": {
        "IM_EXHIBIT_TITLE": "Revenue is projected to triple within 24 months",
        "IM_EXHIBIT_TAKEAWAY": "Growth is driven primarily by Segment B, which the pilot validated as the highest-margin entry point.",
        "IM_EXHIBIT_SOURCE": "Source: internal model, client-provided FY25 actuals",
    },
    "quote_callout": {
        "IM_QUOTE_TEXT": "This is the clearest market-entry case we have reviewed in three years — the numbers hold up under every stress test.",
        "IM_QUOTE_ATTRIB": "— Investment Committee Chair, client organization",
    },
    "closing_contact": {
        "IM_CLOSING_CONTACT": "Prit Dhanani\nAI Engineer\nprit.dhanani@example.com · +1 (555) 010-1234",
    },
}


def build():
    template_prs = Presentation(str(TEMPLATE_PATH))
    manifest = load_manifest()
    n_seed_slides = len(template_prs.slides)  # 9, before we add anything

    for layout in manifest["layouts"]:
        spec = SAMPLE_SPECS[layout["layout_id"]]
        images = {name: path for name, path in IMAGE_FOR_SLOT.items()
                  if any(s["shape_name"] == name for s in layout["slots"])}
        render_slide(template_prs, template_prs, layout, spec, image_paths=images)

    # Drop the original 9 seed/index slides, keep only the 8 newly rendered ones.
    id_lst = template_prs.slides._sldIdLst
    sld_ids = list(id_lst)
    for i in range(n_seed_slides - 1, -1, -1):
        id_lst.remove(sld_ids[i])

    template_prs.save(str(OUT_PATH))
    print(f"wrote {OUT_PATH} with {len(template_prs.slides)} slides")
    return OUT_PATH


if __name__ == "__main__":
    build()
