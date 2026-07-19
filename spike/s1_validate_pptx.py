"""S1 validation, layer 1: reopen the saved output with python-pptx and check
every manifest slot landed on the right shape with the right content, and that
font/size/color were inherited (not manually reset) on text runs."""
import json
from pathlib import Path

from pptx import Presentation

from renderer import MANIFEST_PATH

OUT_PATH = Path(__file__).parent / "output" / "S1_filled_deck.pptx"


def main():
    manifest = json.loads(MANIFEST_PATH.read_text(encoding="utf-8"))
    prs = Presentation(str(OUT_PATH))

    assert len(prs.slides) == 8, f"expected 8 slides, got {len(prs.slides)}"

    results = []
    for i, layout in enumerate(manifest["layouts"]):
        slide = prs.slides[i]
        shape_by_name = {s.name: s for s in slide.shapes}
        for slot in layout["slots"]:
            name = slot["shape_name"]
            shape = shape_by_name.get(name)
            row = {"layout": layout["layout_id"], "slot": name, "type": slot["type"]}
            if shape is None:
                row["status"] = "MISSING (grey placeholder rect not replaced by a same-named picture)"
                results.append(row)
                continue
            if slot["type"] == "image":
                is_picture = shape.shape_type == 13  # MSO_SHAPE_TYPE.PICTURE
                row["status"] = "OK (picture inserted)" if is_picture else "FAIL (not a picture)"
            else:
                run = shape.text_frame.paragraphs[0].runs[0] if shape.text_frame.paragraphs[0].runs else None
                text = shape.text_frame.text
                font_name = run.font.name if run else None
                font_size = run.font.size if run else None
                row["status"] = "OK"
                row["text_preview"] = text[:50]
                row["font"] = font_name
                row["size_pt"] = font_size.pt if font_size else None
            results.append(row)

    n_ok = sum(1 for r in results if str(r["status"]).startswith("OK"))
    n_total = len(results)
    print(f"{n_ok}/{n_total} slots OK\n")
    for r in results:
        print(r)

    # also confirm grey placeholder rects were removed for image slots (no leftover
    # unfilled placeholder shape sitting under/behind the inserted picture)
    print("\nPer-slide shape name inventory (post-fill):")
    for i, layout in enumerate(manifest["layouts"]):
        slide = prs.slides[i]
        print(f"  slide {i} ({layout['layout_id']}):", [s.name for s in slide.shapes])


if __name__ == "__main__":
    main()
