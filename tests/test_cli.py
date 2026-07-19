import json
import tempfile
import unittest
from pathlib import Path
from unittest import mock

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from deck_engine import cli
from deck_engine.models.slide_intent import SlideIntent
from deck_engine.models.slide_spec import SlideSpec
from deck_engine.template_registry import load_tenant_assets


class TestCliBookendHandling(unittest.TestCase):
    """Regression coverage for a real bug hit during the first successful M2
    end-to-end run: when the title slide's content failed business-rule
    validation on both the initial attempt and the bounded retry (D6), the
    CLI silently dropped the title slide entirely instead of rendering it --
    contradicting the fixed, code-guaranteed bookend guarantee (SDD v1.3/v1.4).
    """

    def setUp(self):
        self.manifest = load_tenant_assets("default").manifest
        self.text_bullets_layout = self.manifest.layout_by_id("text_bullets")

    def _run_with_title_generation_failing_both_attempts(self):
        failure_slots = {"IM_TITLE_DECKTITLE": "A Title", "IM_TITLE_SUBTITLE": "x" * 111}

        def fake_generate_single_slide(brief, layout, provider, model, role="opening slide"):
            if layout.layout_id == cli.TITLE_LAYOUT_ID:
                return None, {"slots": failure_slots, "errors": ["slot 'IM_TITLE_SUBTITLE' exceeds max_chars=110 (got 111)"]}
            return (
                SlideSpec(layout_id=layout.layout_id, content_kind="text", slots={"IM_BODY_BULLETS": ["A point"]}),
                None,
            )

        intent = SlideIntent(
            purpose="Make a point",
            suggested_layout="text_bullets",
            headline="A Headline",
            content_outline="Some outline",
            needs_exhibit="none",
        )

        with tempfile.TemporaryDirectory() as tmpdir:
            out_path = Path(tmpdir) / "deck.pptx"
            with mock.patch("deck_engine.cli.check_provider_ready"), mock.patch(
                "deck_engine.cli.generate_single_slide", side_effect=fake_generate_single_slide
            ), mock.patch("deck_engine.cli.plan_deck", return_value=([intent], None)), mock.patch(
                "deck_engine.cli.select_layout", return_value=(self.text_bullets_layout, None)
            ):
                exit_code = cli.main(
                    [
                        "--brief", "irrelevant, generation is mocked",
                        "--contact", "Jane Doe|Role|jane@example.com",
                        "--out", str(out_path),
                    ]
                )
            self.assertEqual(exit_code, 0)
            prs = Presentation(str(out_path))
            return prs

    def test_title_slide_is_still_rendered_when_both_generation_attempts_fail_validation(self):
        prs = self._run_with_title_generation_failing_both_attempts()
        # title (best-effort, flagged) + 1 content slide + closing = 3 slides.
        # The bug this guards against rendered only 2 (content + closing).
        self.assertEqual(len(prs.slides), 3)
        first_slide_text = " ".join(
            shape.text_frame.text for shape in prs.slides[0].shapes if shape.has_text_frame
        )
        self.assertIn("A Title", first_slide_text)


class TestCliChartMapping(unittest.TestCase):
    """M3: --chart-data blocks map, in planner order, to slides the Planner
    marked needs_exhibit=chart; the Chart Builder resolves a supported type and
    the Renderer produces a native chart. No LLM is exercised (all mocked) --
    this pins the wiring/mapping, not model behavior."""

    def setUp(self):
        self.manifest = load_tenant_assets("default").manifest
        self.exhibit_layout = self.manifest.layout_by_id("exhibit_data")
        self.bullets_layout = self.manifest.layout_by_id("text_bullets")

    def _intents(self):
        return [
            SlideIntent(purpose="Show revenue", suggested_layout="exhibit_data",
                        headline="Revenue", content_outline="Q1-Q4 growth", needs_exhibit="chart"),
            SlideIntent(purpose="Make a point", suggested_layout="text_bullets",
                        headline="Point", content_outline="some text", needs_exhibit="none"),
        ]

    def _fake_generate(self, brief, layout, provider, model, role="opening slide"):
        return SlideSpec(layout_id=layout.layout_id, content_kind="text", slots={}), None

    def _fake_select(self, intent, manifest, provider, model):
        return manifest.layout_by_id(intent.suggested_layout), None

    def _run(self, chart_blocks, extra_args=None):
        with tempfile.TemporaryDirectory() as tmpdir:
            out_path = Path(tmpdir) / "deck.pptx"
            data_path = Path(tmpdir) / "charts.json"
            data_path.write_text(json.dumps(chart_blocks), encoding="utf-8")
            argv = [
                "--brief", "mocked",
                "--contact", "Jane|Role|jane@x.com",
                "--chart-data", str(data_path),
                "--out", str(out_path),
            ] + (extra_args or [])
            with mock.patch("deck_engine.cli.check_provider_ready"), mock.patch(
                "deck_engine.cli.generate_single_slide", side_effect=self._fake_generate
            ), mock.patch(
                "deck_engine.cli.plan_deck", return_value=(self._intents(), None)
            ), mock.patch("deck_engine.cli.select_layout", side_effect=self._fake_select):
                exit_code = cli.main(argv)
            self.assertEqual(exit_code, 0)
            return Presentation(str(out_path))

    def test_chart_block_becomes_a_native_chart_on_the_exhibit_slide(self):
        prs = self._run([
            {"categories": ["Q1", "Q2", "Q3", "Q4"], "series": {"Revenue": [2.1, 2.5, 2.9, 3.4]}, "title": "Rev"}
        ])
        # title + exhibit(chart) + text_bullets + closing = 4 slides
        self.assertEqual(len(prs.slides), 4)
        all_charts = [s for slide in prs.slides for s in slide.shapes if s.has_chart]
        self.assertEqual(len(all_charts), 1)  # exactly one native chart, on the exhibit slide

    def test_type_override_in_block_is_respected(self):
        prs = self._run([
            {"categories": ["Q1", "Q2"], "series": {"r": [1, 2]}, "type": "bar"}
        ])
        charts = [s for slide in prs.slides for s in slide.shapes if s.has_chart]
        self.assertEqual(len(charts), 1)
        self.assertEqual(str(charts[0].chart.chart_type), "BAR_CLUSTERED (57)")

    def test_chart_needing_slide_without_data_keeps_placeholder_flagged(self):
        # no chart blocks at all -> exhibit slot stays a placeholder, not a crash
        prs = self._run([])
        charts = [s for slide in prs.slides for s in slide.shapes if s.has_chart]
        self.assertEqual(len(charts), 0)
        self.assertEqual(len(prs.slides), 4)  # still rendered, just without the chart


class TestCliDiagramMapping(unittest.TestCase):
    """M4: --diagram-data blocks map, in planner order, to slides the Planner
    marked needs_exhibit=diagram; the Diagram Builder validates the steps and
    the Renderer produces a native grouped chevron-flow. No LLM is exercised
    (all mocked) -- this pins the wiring/mapping, not model behavior."""

    def setUp(self):
        self.manifest = load_tenant_assets("default").manifest

    def _intents(self):
        return [
            SlideIntent(purpose="Show the process", suggested_layout="image_only",
                        headline="How it works", content_outline="the flow", needs_exhibit="diagram"),
            SlideIntent(purpose="Make a point", suggested_layout="text_bullets",
                        headline="Point", content_outline="some text", needs_exhibit="none"),
        ]

    def _fake_generate(self, brief, layout, provider, model, role="opening slide"):
        return SlideSpec(layout_id=layout.layout_id, content_kind="text", slots={}), None

    def _fake_select(self, intent, manifest, provider, model):
        return manifest.layout_by_id(intent.suggested_layout), None

    def _run(self, diagram_blocks, extra_args=None):
        with tempfile.TemporaryDirectory() as tmpdir:
            out_path = Path(tmpdir) / "deck.pptx"
            data_path = Path(tmpdir) / "diagrams.json"
            data_path.write_text(json.dumps(diagram_blocks), encoding="utf-8")
            argv = [
                "--brief", "mocked",
                "--contact", "Jane|Role|jane@x.com",
                "--diagram-data", str(data_path),
                "--out", str(out_path),
            ] + (extra_args or [])
            with mock.patch("deck_engine.cli.check_provider_ready"), mock.patch(
                "deck_engine.cli.generate_single_slide", side_effect=self._fake_generate
            ), mock.patch(
                "deck_engine.cli.plan_deck", return_value=(self._intents(), None)
            ), mock.patch("deck_engine.cli.select_layout", side_effect=self._fake_select):
                exit_code = cli.main(argv)
            self.assertEqual(exit_code, 0)
            return Presentation(str(out_path))

    def _groups(self, prs):
        return [s for slide in prs.slides for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.GROUP]

    def test_diagram_block_becomes_a_grouped_chevron_flow_on_the_exhibit_slide(self):
        prs = self._run([{"steps": ["Intake", "Review", "Deliver"]}])
        # title + image_only(diagram) + text_bullets + closing = 4 slides
        self.assertEqual(len(prs.slides), 4)
        groups = self._groups(prs)
        self.assertEqual(len(groups), 1)  # exactly one chevron-flow group
        self.assertEqual(len(groups[0].shapes), 3)  # one chevron per step

    def test_diagram_needing_slide_without_data_keeps_placeholder_flagged(self):
        # no diagram blocks at all -> hero slot stays a placeholder, not a crash
        prs = self._run([])
        self.assertEqual(len(self._groups(prs)), 0)
        self.assertEqual(len(prs.slides), 4)  # still rendered, just without the diagram

    def test_invalid_diagram_block_is_flagged_and_skipped_not_crashed(self):
        # 1 step violates MIN_STEPS -> block is rejected, placeholder stays, run still succeeds
        prs = self._run([{"steps": ["only one step"]}])
        self.assertEqual(len(self._groups(prs)), 0)
        self.assertEqual(len(prs.slides), 4)


class TestSecondTenantCli(unittest.TestCase):
    """M7 (SDD v1.9): the exact same CLI, unmodified, generates a real deck
    against the second tenant's own template/manifest/brand -- proving
    tenant-agnosticism end-to-end (PRD FR-11/FR-12), not just that the
    manifest loads."""

    def setUp(self):
        self.manifest = load_tenant_assets("meridian").manifest

    def _intents(self):
        return [
            SlideIntent(purpose="Show the process", suggested_layout="image_only",
                        headline="How it works", content_outline="the flow", needs_exhibit="diagram"),
        ]

    def _fake_generate(self, brief, layout, provider, model, role="opening slide"):
        return SlideSpec(layout_id=layout.layout_id, content_kind="text", slots={}), None

    def _fake_select(self, intent, manifest, provider, model):
        return manifest.layout_by_id(intent.suggested_layout), None

    def test_generates_a_valid_deck_with_meridians_own_brand_color_and_shape_names(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            out_path = Path(tmpdir) / "meridian_deck.pptx"
            data_path = Path(tmpdir) / "steps.json"
            data_path.write_text(json.dumps([{"steps": ["Discover", "Design", "Deliver"]}]), encoding="utf-8")
            argv = [
                "--tenant", "meridian",
                "--brief", "mocked", "--contact", "Jane|Role|jane@x.com",
                "--diagram-data", str(data_path),
                "--out", str(out_path),
            ]
            with mock.patch("deck_engine.cli.check_provider_ready"), mock.patch(
                "deck_engine.cli.generate_single_slide", side_effect=self._fake_generate
            ), mock.patch(
                "deck_engine.cli.plan_deck", return_value=(self._intents(), None)
            ), mock.patch("deck_engine.cli.select_layout", side_effect=self._fake_select):
                exit_code = cli.main(argv)
            self.assertEqual(exit_code, 0)
            prs = Presentation(str(out_path))

        groups = [s for slide in prs.slides for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.GROUP]
        self.assertEqual(len(groups), 1)
        chevron = groups[0].shapes[0]
        self.assertEqual(str(chevron.fill.fore_color.rgb), self.manifest.brand["colors"]["teal_primary"])
        all_shape_names = {s.name for slide in prs.slides for s in slide.shapes}
        self.assertTrue(any(name.startswith("MP_") for name in all_shape_names))


if __name__ == "__main__":
    unittest.main()
