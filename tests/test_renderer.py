import unittest

from pptx.enum.shapes import MSO_SHAPE_TYPE

from deck_engine.chart_builder import build_chart_spec
from deck_engine.diagram_builder import build_diagram_data
from deck_engine.models.chart import ChartData
from deck_engine.models.manifest import GeometryIn, Layout, Slot
from deck_engine.models.slide_spec import SlideSpec
from deck_engine.renderer import fill_slot, find_shape, render_slide, replace_rendered_slide, strip_seed_slides
from deck_engine.template_registry import load_tenant_assets


class TestRenderer(unittest.TestCase):
    def setUp(self):
        self.assets = load_tenant_assets("default")
        self.layout = self.assets.manifest.layout_by_id("text_bullets")

    def test_render_slide_fills_text_and_bullets_correctly(self):
        prs = self.assets.open_template()
        n_seed_slides = len(prs.slides)
        seed_slide = prs.slides[self.layout.slide_index]
        original_seed_title_text = find_shape(seed_slide, "IM_BULLETS_TITLE").text_frame.text

        spec = SlideSpec(
            layout_id="text_bullets",
            slots={
                "IM_BULLETS_TITLE": "Test Title",
                "IM_BULLETS_KICKER": "Test Kicker",
                "IM_BULLETS_BODY": ["first point", "second point", "third point"],
            },
        )
        new_slide, skipped, _ = render_slide(prs, seed_slide, self.layout, spec)

        self.assertEqual(skipped, [])
        self.assertEqual(find_shape(new_slide, "IM_BULLETS_TITLE").text_frame.text, "Test Title")
        self.assertEqual(find_shape(new_slide, "IM_BULLETS_KICKER").text_frame.text, "Test Kicker")
        body_text = find_shape(new_slide, "IM_BULLETS_BODY").text_frame.text
        self.assertIn("first point", body_text)
        self.assertIn("second point", body_text)
        self.assertIn("third point", body_text)

        # D3/D13: duplication must not mutate the pristine seed slide itself.
        self.assertEqual(
            find_shape(seed_slide, "IM_BULLETS_TITLE").text_frame.text,
            original_seed_title_text,
        )
        self.assertNotEqual(new_slide, seed_slide)

    def test_render_slide_inherits_font_from_seed(self):
        prs = self.assets.open_template()
        seed_slide = prs.slides[self.layout.slide_index]
        spec = SlideSpec(layout_id="text_bullets", slots={"IM_BULLETS_TITLE": "Font Check", "IM_BULLETS_BODY": ["a"]})
        new_slide, _, _ = render_slide(prs, seed_slide, self.layout, spec)

        run = find_shape(new_slide, "IM_BULLETS_TITLE").text_frame.paragraphs[0].runs[0]
        self.assertEqual(run.font.name, "Cambria")

    def test_strip_seed_slides_leaves_only_rendered_slides(self):
        prs = self.assets.open_template()
        n_seed_slides = len(prs.slides)
        seed_slide = prs.slides[self.layout.slide_index]

        spec = SlideSpec(layout_id="text_bullets", slots={"IM_BULLETS_TITLE": "Only Slide", "IM_BULLETS_BODY": ["a"]})
        render_slide(prs, seed_slide, self.layout, spec)
        strip_seed_slides(prs, n_seed_slides)

        self.assertEqual(len(prs.slides), 1)
        self.assertEqual(find_shape(prs.slides[0], "IM_BULLETS_TITLE").text_frame.text, "Only Slide")

    def test_multiple_renders_always_duplicate_the_same_pristine_seed(self):
        prs = self.assets.open_template()
        n_seed_slides = len(prs.slides)
        seed_slide = prs.slides[self.layout.slide_index]

        for i in range(3):
            spec = SlideSpec(layout_id="text_bullets", slots={"IM_BULLETS_TITLE": f"Slide {i}", "IM_BULLETS_BODY": ["a"]})
            render_slide(prs, seed_slide, self.layout, spec)  # always the same seed_slide object (D13)

        strip_seed_slides(prs, n_seed_slides)
        self.assertEqual(len(prs.slides), 3)
        titles = [find_shape(s, "IM_BULLETS_TITLE").text_frame.text for s in prs.slides]
        self.assertEqual(titles, ["Slide 0", "Slide 1", "Slide 2"])

    def test_missing_required_image_is_skipped_not_crashed_and_flagged(self):
        """M2: layouts needing an exhibit (two_column, image_only, exhibit_data)
        get selected before the Chart/Diagram Builders (M3/M4) exist to supply
        an image. render_slide must not crash -- it leaves the seed's grey
        placeholder in place and reports the slot back so the caller can
        flag needs_review (D9), rather than shipping a silently-broken slide.
        """
        two_column = self.assets.manifest.layout_by_id("two_column")
        prs = self.assets.open_template()
        seed_slide = prs.slides[two_column.slide_index]
        spec = SlideSpec(
            layout_id="two_column",
            slots={"IM_TWOCOL_TITLE": "No Image Available", "IM_TWOCOL_BODY": ["a", "b"]},
        )
        new_slide, skipped, _ = render_slide(prs, seed_slide, two_column, spec)  # no image_paths given

        self.assertEqual(skipped, ["IM_TWOCOL_IMAGE"])
        # the grey placeholder shape is still there, still named as the manifest expects
        self.assertIsNotNone(find_shape(new_slide, "IM_TWOCOL_IMAGE"))

    def test_chart_spec_fills_exhibit_slot_with_native_editable_chart(self):
        """M3: an exhibit image slot given a ChartSpec renders a real, native
        python-pptx chart (editable, not a picture) at the slot's geometry,
        named back to the manifest shape_name, with the grey placeholder rect
        removed -- no required slot left skipped (5.8 point 5 / D7)."""
        exhibit = self.assets.manifest.layout_by_id("exhibit_data")
        img_slot = exhibit.image_slots()[0]
        prs = self.assets.open_template()
        seed_slide = prs.slides[exhibit.slide_index]

        data = ChartData(categories=["Q1", "Q2", "Q3", "Q4"], series={"Revenue": [2.1, 2.5, 2.9, 3.4]}, title="Rev")
        chart_spec = build_chart_spec(data)  # single time series -> line
        spec = SlideSpec(
            layout_id="exhibit_data",
            content_kind="chart",
            slots={"IM_EXHIBIT_TITLE": "Growth", "IM_EXHIBIT_TAKEAWAY": "Up and to the right."},
        )
        new_slide, skipped, _ = render_slide(
            prs, seed_slide, exhibit, spec, chart_specs={img_slot.shape_name: chart_spec}
        )

        self.assertEqual(skipped, [])  # the exhibit slot was filled by the chart, not left skipped
        chart_shapes = [s for s in new_slide.shapes if s.has_chart]
        self.assertEqual(len(chart_shapes), 1)
        self.assertEqual(chart_shapes[0].name, img_slot.shape_name)  # stays addressable (D3)

    def test_chart_takes_precedence_over_image_path_for_exhibit_slot(self):
        """When both a chart spec and an image path are offered for the same
        slot, the native (editable) chart wins (5.8 point 5 calls it preferred)."""
        exhibit = self.assets.manifest.layout_by_id("exhibit_data")
        img_slot = exhibit.image_slots()[0]
        prs = self.assets.open_template()
        seed_slide = prs.slides[exhibit.slide_index]

        data = ChartData(categories=["A", "B"], series={"x": [1, 2]})
        chart_spec = build_chart_spec(data)
        spec = SlideSpec(
            layout_id="exhibit_data",
            content_kind="chart",
            slots={"IM_EXHIBIT_TITLE": "T", "IM_EXHIBIT_TAKEAWAY": "point"},  # fill the required text slots
        )
        new_slide, skipped, _ = render_slide(
            prs, seed_slide, exhibit, spec,
            image_paths={img_slot.shape_name: "does_not_exist.png"},  # would fail if the picture path were used
            chart_specs={img_slot.shape_name: chart_spec},
        )
        self.assertEqual(skipped, [])
        self.assertEqual(len([s for s in new_slide.shapes if s.has_chart]), 1)

    def test_diagram_data_fills_exhibit_slot_with_grouped_chevron_shapes(self):
        """M4: an exhibit image slot given a DiagramData renders a native
        chevron-flow diagram -- ONE group shape named back to the manifest
        shape_name (D3), whose children are the individually-editable chevrons
        in step order, with the grey placeholder rect removed (5.7)."""
        image_only = self.assets.manifest.layout_by_id("image_only")
        img_slot = image_only.image_slots()[0]
        prs = self.assets.open_template()
        seed_slide = prs.slides[image_only.slide_index]

        labels = ["Intake", "Review", "Approve", "Deliver"]
        diagram_data = build_diagram_data(labels)
        spec = SlideSpec(layout_id="image_only", content_kind="diagram", slots={})
        new_slide, skipped, _ = render_slide(
            prs, seed_slide, image_only, spec, diagram_specs={img_slot.shape_name: diagram_data}
        )

        self.assertEqual(skipped, [])  # the hero slot was filled by the diagram
        groups = [s for s in new_slide.shapes if s.shape_type == MSO_SHAPE_TYPE.GROUP]
        self.assertEqual(len(groups), 1)
        self.assertEqual(groups[0].name, img_slot.shape_name)  # stays addressable (D3)
        self.assertEqual(len(groups[0].shapes), len(labels))
        self.assertEqual([c.text_frame.text for c in groups[0].shapes], labels)

    def test_diagram_takes_precedence_over_image_path_for_exhibit_slot(self):
        """When both a diagram and an image path are offered for the same slot,
        the native (editable) diagram wins -- mirroring the chart precedence."""
        image_only = self.assets.manifest.layout_by_id("image_only")
        img_slot = image_only.image_slots()[0]
        prs = self.assets.open_template()
        seed_slide = prs.slides[image_only.slide_index]

        diagram_data = build_diagram_data(["One", "Two"])
        spec = SlideSpec(layout_id="image_only", content_kind="diagram", slots={})
        new_slide, skipped, _ = render_slide(
            prs, seed_slide, image_only, spec,
            image_paths={img_slot.shape_name: "does_not_exist.png"},  # would fail if the picture path were used
            diagram_specs={img_slot.shape_name: diagram_data},
        )
        self.assertEqual(skipped, [])
        self.assertEqual(len([s for s in new_slide.shapes if s.shape_type == MSO_SHAPE_TYPE.GROUP]), 1)

    def test_chart_takes_precedence_over_diagram_for_same_slot(self):
        """Defensive ordering: if both a chart spec and a diagram are supplied
        for the same slot, exactly one native chart is built and no diagram
        group appears (fill_slot checks chart first)."""
        exhibit = self.assets.manifest.layout_by_id("exhibit_data")
        img_slot = exhibit.image_slots()[0]
        prs = self.assets.open_template()
        seed_slide = prs.slides[exhibit.slide_index]

        chart_spec = build_chart_spec(ChartData(categories=["A", "B"], series={"x": [1, 2]}))
        diagram_data = build_diagram_data(["One", "Two"])
        spec = SlideSpec(
            layout_id="exhibit_data",
            content_kind="chart",
            slots={"IM_EXHIBIT_TITLE": "T", "IM_EXHIBIT_TAKEAWAY": "point"},
        )
        new_slide, skipped, _ = render_slide(
            prs, seed_slide, exhibit, spec,
            chart_specs={img_slot.shape_name: chart_spec},
            diagram_specs={img_slot.shape_name: diagram_data},
        )
        self.assertEqual(skipped, [])
        self.assertEqual(len([s for s in new_slide.shapes if s.has_chart]), 1)
        self.assertEqual(len([s for s in new_slide.shapes if s.shape_type == MSO_SHAPE_TYPE.GROUP]), 0)

    def test_fill_slot_flags_overflow_for_a_tiny_box_with_long_text(self):
        """M5: fill_slot runs the QA/Overflow Checker (5.9) on text/bullets
        slots when given layout+slide_height_in. A real shape/font from the
        template, deliberately given a tiny synthetic geometry (no other
        slots in this layout, so headroom falls back to the slide's bottom
        margin) and a long value, must be flagged -- never silently accepted."""
        prs = self.assets.open_template()
        seed_slide = prs.slides[self.layout.slide_index]
        tiny_slot = Slot(
            shape_name="IM_BULLETS_TITLE", type="text", role="title",
            geometry_in=GeometryIn(left_in=1.0, top_in=1.0, width_in=1.0, height_in=0.1),
        )
        tiny_layout = Layout(layout_id="tiny", slide_index=self.layout.slide_index, title="tiny", use_when="x", slots=[tiny_slot])
        long_text = "This is a long line of text that will not fit inside a one inch wide, tenth-of-an-inch tall box."

        filled, overflow_reason = fill_slot(seed_slide, tiny_slot, long_text, layout=tiny_layout, slide_height_in=7.5)

        self.assertTrue(filled)
        self.assertIsNotNone(overflow_reason)
        self.assertIn("overflow", overflow_reason)

    def test_fill_slot_skips_overflow_check_when_not_opted_in(self):
        """Backward compatibility: omitting layout/slide_height_in (the
        default) skips the overflow check entirely -- existing callers that
        don't care about it are unaffected."""
        prs = self.assets.open_template()
        seed_slide = prs.slides[self.layout.slide_index]
        tiny_slot = Slot(
            shape_name="IM_BULLETS_TITLE", type="text", role="title",
            geometry_in=GeometryIn(left_in=1.0, top_in=1.0, width_in=1.0, height_in=0.1),
        )
        long_text = "This is a long line of text that will not fit inside a one inch wide, tenth-of-an-inch tall box."

        filled, overflow_reason = fill_slot(seed_slide, tiny_slot, long_text)  # no layout/slide_height_in

        self.assertTrue(filled)
        self.assertIsNone(overflow_reason)

    def test_render_slide_reports_no_overflow_for_normal_short_content(self):
        """Regression guard: ordinary, well within-limits content on the real
        manifest's own layout/geometry must not spuriously flag overflow."""
        prs = self.assets.open_template()
        seed_slide = prs.slides[self.layout.slide_index]
        spec = SlideSpec(
            layout_id="text_bullets",
            slots={
                "IM_BULLETS_TITLE": "Test Title",
                "IM_BULLETS_KICKER": "Test Kicker",
                "IM_BULLETS_BODY": ["first point", "second point", "third point"],
            },
        )
        _, _, overflow = render_slide(prs, seed_slide, self.layout, spec, slide_height_in=7.5)
        self.assertEqual(overflow, [])

    def test_replace_rendered_slide_preserves_order_and_discards_the_stale_slide(self):
        """M6/5.10: revising one slide must not disturb any other slide's
        order or content (FR-6), and must re-duplicate the pristine seed
        fresh rather than mutate the already-rendered slide (D13)."""
        prs = self.assets.open_template()
        n_seed = len(prs.slides)
        seed_slide = prs.slides[self.layout.slide_index]

        titles = ["Slide 0", "Slide 1", "Slide 2"]
        for t in titles:
            render_slide(prs, seed_slide, self.layout, SlideSpec(layout_id="text_bullets", slots={"IM_BULLETS_TITLE": t, "IM_BULLETS_BODY": ["a"]}))

        revised_spec = SlideSpec(layout_id="text_bullets", slots={"IM_BULLETS_TITLE": "Slide 1 REVISED", "IM_BULLETS_BODY": ["a"]})
        position = n_seed + 1  # "Slide 1"'s position among the rendered slides
        new_slide, skipped, _ = replace_rendered_slide(prs, position, seed_slide, self.layout, revised_spec)

        self.assertEqual(skipped, [])
        rendered_titles = [find_shape(prs.slides[n_seed + i], "IM_BULLETS_TITLE").text_frame.text for i in range(3)]
        self.assertEqual(rendered_titles, ["Slide 0", "Slide 1 REVISED", "Slide 2"])
        self.assertEqual(len(prs.slides), n_seed + 3)  # no extra slide left behind
        self.assertEqual(prs.slides[n_seed + 1], new_slide)  # the slide at that position IS the new one


if __name__ == "__main__":
    unittest.main()
