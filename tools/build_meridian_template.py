"""One-time template-authoring tool for M7's second tenant ("Meridian
Partners"), proving tenant-agnosticism (SDD §9 M7, PRD FR-11/FR-12).

A real onboarding builds the seed slides by hand in PowerPoint (per the
manifest's own authoring note: "~1-2 hours per template"). This script is a
programmatic stand-in for that manual step -- there's no design software in
this environment -- but it produces a structurally real artifact: seed slides
on the same generic Blank PowerPoint layout (D3), with uniquely-named shapes
carrying real font/color formatting on their runs (so the renderer's
format-preserving fill, which reads existing run formatting rather than
setting it, genuinely inherits Meridian's brand, not the reference tenant's).

Deliberately differs from MasterDeck in every way that matters for proving
isolation, not just cosmetically:
  - Different shape-name prefix (MP_ vs IM_) and different per-slot names.
  - Different brand colors and fonts (Georgia/Verdana vs Cambria/Calibri).
  - Different geometry (0.75in margins vs 0.9in) -- proves nothing in the
    renderer/chart/diagram code hardcodes the reference tenant's coordinates.
  - Different, deliberately shuffled slide_index per layout_id (quote_callout
    is slide 0, title is slide 1, ...) -- proves layout_id -> slide_index is
    fully manifest-driven, never assumed to follow a fixed physical order.

The single LAYOUTS list below is the one source of truth for both the .pptx
shapes and the manifest JSON, so the two can never drift apart.
"""
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

REPO_ROOT = Path(__file__).resolve().parent.parent
OUT_PPTX = REPO_ROOT / "Templates" / "Meridian.pptx"
OUT_MANIFEST = REPO_ROOT / "Templates" / "Meridian_manifest.json"

SLIDE_W_IN, SLIDE_H_IN = 13.333, 7.5

BRAND = {
    "colors": {
        "primary": "0E6E5C",  # the tenant-agnostic key the engine looks up (e.g. diagram chevrons)
        "teal_primary": "0E6E5C",
        "amber_accent": "E2892A",
        "ink": "1B1B1B",
        "muted": "707070",
        "surface": "FFFFFF",
        "placeholder_fill": "E7E7E7",
    },
    "wordmark": "MERIDIAN PARTNERS",
    "claims": ["Clarity you can act on."],
    "fonts": {
        "display": "Georgia",
        "body": "Verdana",
        "note": "Distinct stand-in fonts (vs the reference tenant's Cambria/Calibri) -- proves font "
        "inheritance is manifest/template-driven, not hardcoded anywhere in the renderer (M7).",
    },
}
TITLE_FONT_SIZE = 28
BODY_FONT_SIZE = 14
PLACEHOLDER_FONT_SIZE = 13

# (layout_id, slide_index, title, use_when, [slot dicts]) -- slide_index is
# deliberately NOT in layout-definition order, to prove nothing assumes one.
LAYOUTS = [
    ("quote_callout", 0, "Quote / Callout", "A pull quote, client testimonial, or single big statement.", [
        {"slot": "STATEMENT", "type": "text", "role": "quote", "required": True, "max_chars": 160,
         "geom": (1.4, 2.05, 10.53, 2.85)},
        {"slot": "BYLINE", "type": "text", "role": "attribution", "required": False, "max_chars": 70,
         "geom": (1.42, 5.0, 10.0, 0.6)},
    ]),
    ("title", 1, "Title Slide", "Opening slide of any deck.", [
        {"slot": "HEADLINE", "type": "text", "role": "deck_title", "required": True, "max_chars": 70,
         "geom": (0.75, 4.0, 11.83, 1.2)},
        {"slot": "SUBHEAD", "type": "text", "role": "subtitle", "required": False, "max_chars": 110,
         "geom": (0.77, 5.2, 11.4, 0.7)},
        {"slot": "TAGLINE", "type": "text", "role": "claim", "required": False,
         "allowed_values": ["Clarity you can act on."], "geom": (0.77, 2.6, 8.2, 0.5)},
        {"slot": "METALINE", "type": "text", "role": "meta_line", "required": False, "max_chars": 90,
         "geom": (0.77, 6.65, 11.4, 0.5)},
    ]),
    ("closing_contact", 2, "Closing / Contact", "Final slide -- wordmark, claim, contact block.", [
        {"slot": "DETAILS", "type": "text", "role": "contact_block", "required": True, "max_chars": 220,
         "note": "Newline-separated lines: name / role / email x phone.", "geom": (0.77, 4.65, 11.4, 1.6)},
    ]),
    ("text_bullets", 3, "Text + Bullets (no image)", "Pure argument / list slide with no visual.", [
        {"slot": "HEADLINE", "type": "text", "role": "title", "required": True, "max_chars": 60,
         "geom": (0.75, 1.1, 11.83, 0.9)},
        {"slot": "KICKER", "type": "text", "role": "kicker", "required": False, "max_chars": 90,
         "geom": (0.77, 2.0, 11.83, 0.5)},
        {"slot": "LIST", "type": "bullets", "role": "body", "required": True, "max_items": 6,
         "max_chars_per_item": 95, "geom": (0.77, 2.7, 11.83, 3.95)},
    ]),
    ("exhibit_data", 4, "Exhibit / Data", "A chart or table with a so-what takeaway and source.", [
        {"slot": "HEADLINE", "type": "text", "role": "title", "required": True, "max_chars": 60,
         "geom": (0.75, 1.1, 11.83, 0.85)},
        {"slot": "VISUAL", "type": "image", "role": "exhibit", "required": True, "fit": "contain",
         "note": "Insert chart/table image, or overlay a native chart at this geometry.",
         "geom": (0.77, 2.05, 8.6, 4.4)},
        {"slot": "TAKEAWAY", "type": "text", "role": "takeaway", "required": True, "max_chars": 180,
         "geom": (9.55, 2.05, 2.78, 4.05)},
        {"slot": "SOURCE", "type": "text", "role": "source", "required": False, "max_chars": 120,
         "geom": (0.77, 6.5, 9.2, 0.4)},
    ]),
    ("image_only", 5, "Image Only", "A single dominant visual -- exhibit photo, diagram, screenshot.", [
        {"slot": "HERO", "type": "image", "role": "hero_image", "required": True, "fit": "cover",
         "note": "Replace this placeholder rectangle; use its geometry.", "geom": (0.75, 0.85, 11.83, 5.45)},
        {"slot": "CAPTION", "type": "text", "role": "caption", "required": False, "max_chars": 110,
         "geom": (0.77, 6.4, 9.2, 0.5)},
    ]),
    ("section_divider", 6, "Section Divider", "Break between major sections of the deck.", [
        {"slot": "NUM", "type": "text", "role": "section_number", "required": True, "max_chars": 3,
         "geom": (0.75, 2.0, 3.0, 1.4)},
        {"slot": "HEADLINE", "type": "text", "role": "section_title", "required": True, "max_chars": 48,
         "geom": (0.77, 3.35, 11.83, 1.6)},
        {"slot": "KICKER", "type": "text", "role": "kicker", "required": False, "max_chars": 90,
         "geom": (0.8, 5.0, 11.4, 0.6)},
    ]),
    ("two_column", 7, "Two-Column (text + image)", "Narrative on the left, supporting visual on the right.", [
        {"slot": "HEADLINE", "type": "text", "role": "title", "required": True, "max_chars": 60,
         "geom": (0.75, 1.1, 11.83, 0.9)},
        {"slot": "BODY", "type": "bullets", "role": "body", "required": True, "max_items": 5,
         "max_chars_per_item": 80, "geom": (0.77, 2.3, 5.85, 4.15)},
        {"slot": "VISUAL", "type": "image", "role": "support_image", "required": True, "fit": "cover",
         "geom": (6.8, 2.1, 5.78, 4.4)},
    ]),
]


def shape_name(layout_id, slot):
    return f"MP_{layout_id.upper()}_{slot}"


def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank_layout = prs.slide_layouts[6]

    by_index = sorted(LAYOUTS, key=lambda entry: entry[1])
    for layout_id, slide_index, _title, _use_when, slots in by_index:
        slide = prs.slides.add_slide(blank_layout)
        assert len(prs.slides) - 1 == slide_index, "LAYOUTS slide_index must match append order"
        for slot_def in slots:
            left, top, width, height = slot_def["geom"]
            name = shape_name(layout_id, slot_def["slot"])
            if slot_def["type"] == "image":
                shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
                shp.name = name
                shp.fill.solid()
                shp.fill.fore_color.rgb = RGBColor.from_string(BRAND["colors"]["placeholder_fill"])
                shp.line.color.rgb = RGBColor.from_string(BRAND["colors"]["muted"])
                shp.text_frame.text = slot_def["role"].replace("_", " ").upper()
                run = shp.text_frame.paragraphs[0].runs[0]
                run.font.size = Pt(PLACEHOLDER_FONT_SIZE)
                run.font.color.rgb = RGBColor.from_string(BRAND["colors"]["muted"])
                run.font.name = BRAND["fonts"]["body"]
                shp.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                continue

            shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
            shp.name = name
            shp.fill.background()
            shp.line.fill.background()
            is_display = slot_def["role"] in ("deck_title", "title", "section_title", "quote", "claim")
            font_name = BRAND["fonts"]["display"] if is_display else BRAND["fonts"]["body"]
            font_size = TITLE_FONT_SIZE if is_display else BODY_FONT_SIZE
            font_color = BRAND["colors"]["teal_primary"] if is_display else BRAND["colors"]["ink"]

            placeholder_text = f"[{slot_def['role'].replace('_', ' ').title()} placeholder]"
            tf = shp.text_frame
            tf.word_wrap = True
            if slot_def["type"] == "bullets":
                items = ["First point placeholder", "Second point placeholder"]
                tf.text = items[0]
                for item in items[1:]:
                    p = tf.add_paragraph()
                    p.text = item
                paragraphs = tf.paragraphs
            else:
                tf.text = placeholder_text
                paragraphs = [tf.paragraphs[0]]
            for p in paragraphs:
                for run in p.runs:
                    run.font.name = font_name
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = RGBColor.from_string(font_color)

    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPTX))
    print(f"wrote {OUT_PPTX} ({len(prs.slides)} slides)")


def build_manifest():
    manifest = {
        "template_name": "Meridian Partners Deck",
        "template_version": "1.0",
        "pipeline": "Claude -> JSON slide-specs -> python-pptx renderer",
        "fill_protocol": [
            "1. Pick a layout by 'layout_id'.",
            "2. Duplicate the matching template slide (slide_index).",
            "3. For each slot, find the shape where shape.name == slot.shape_name.",
            "4. text  -> set run text on that shape.",
            "5. bullets -> one paragraph per item (respect max_items / max_chars_per_item).",
            "6. image -> insert a picture at slot.geometry_in (inches); remove the grey placeholder rect.",
            "7. Never add unnamed shapes; never rename existing ones.",
        ],
        "brand": BRAND,
        "slide_dimensions_in": {"width": SLIDE_W_IN, "height": SLIDE_H_IN},
        "shape_name_convention": "MP_<LAYOUT>_<SLOT>  (uppercase; unique per slide)",
        "layouts": [],
    }
    for layout_id, slide_index, title, use_when, slots in LAYOUTS:
        entry = {"layout_id": layout_id, "slide_index": slide_index, "title": title, "use_when": use_when, "slots": []}
        for slot_def in slots:
            left, top, width, height = slot_def["geom"]
            slot = {
                "shape_name": shape_name(layout_id, slot_def["slot"]),
                "geometry_in": {"left_in": left, "top_in": top, "width_in": width, "height_in": height},
                "type": slot_def["type"],
                "role": slot_def["role"],
                "required": slot_def["required"],
            }
            for optional_key in ("max_chars", "max_items", "max_chars_per_item", "allowed_values", "fit", "note"):
                if optional_key in slot_def:
                    slot[optional_key] = slot_def[optional_key]
            entry["slots"].append(slot)
        manifest["layouts"].append(entry)

    OUT_MANIFEST.write_text(json.dumps(manifest, indent=2) + "\n", encoding="utf-8")
    print(f"wrote {OUT_MANIFEST} ({len(manifest['layouts'])} layouts)")


if __name__ == "__main__":
    build_pptx()
    build_manifest()
